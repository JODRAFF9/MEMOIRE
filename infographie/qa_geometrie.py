"""QA geometrique : marges, chevauchements, debordement estime (Calibri).
Parcourt toutes les diapositives du fichier passe en argument."""
from pptx import Presentation
import sys
EMU = 914400.0
f = sys.argv[1] if len(sys.argv) > 1 else 'infographie_memoire.pptx'
p = Presentation(f); W, H = p.slide_width/EMU, p.slide_height/EMU
print(f"Diapositive {W:.2f} x {H:.2f} pouces - {len(p.slides)} diapositives")

def infos(sh):
    sz, ls, bold = 12, None, False
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            if para.line_spacing is not None:
                ls = para.line_spacing.pt if hasattr(para.line_spacing, 'pt') else float(para.line_spacing)*sz
            for r in para.runs:
                if r.font.size: sz = r.font.size.pt
                if r.font.bold: bold = True
    return sz, (ls or sz*1.25), bold

total = 0
for idx, slide in enumerate(p.slides, 1):
    boites = []
    for sh in slide.shapes:
        t = sh.text_frame.text if sh.has_text_frame else ""
        sz, ls, bold = infos(sh)
        boites.append(dict(n=sh.name, x=sh.left/EMU, y=sh.top/EMU, w=sh.width/EMU,
                           h=sh.height/EMU, t=t, sz=sz, ls=ls, bold=bold))
    msgs = []

    for b in boites:
        if not (b['t'] or 'Chart' in b['n'] or 'Image' in b['n']): continue
        if b['x'] < 0.40 or b['y'] < 0.15 or b['x']+b['w'] > W-0.40 or b['y']+b['h'] > H-0.12:
            msgs.append(f"  MARGE    {b['n']} [{b['x']:.2f},{b['y']:.2f} {b['w']:.2f}x{b['h']:.2f}] «{b['t'][:35]}»")

    for b in boites:
        if not b['t']: continue
        car = b['sz'] * (0.50 if b['bold'] else 0.475) / 72
        cpl = max(1, int(b['w'] / car))
        lignes = sum(max(1, -(-len(l)//cpl)) for l in b['t'].split("\n"))
        besoin = lignes * b['ls'] / 72
        if besoin > b['h'] + 0.04:
            msgs.append(f"  DEBORD   {b['n']} {b['sz']:.0f}pt {lignes}L besoin {besoin:.2f}\"/{b['h']:.2f}\" «{b['t'][:45]}»")

    vus = set()
    for i, a in enumerate(boites):
        for c in boites[i+1:]:
            if not (a['t'] and c['t']): continue
            ox = min(a['x']+a['w'], c['x']+c['w']) - max(a['x'], c['x'])
            oy = min(a['y']+a['h'], c['y']+c['h']) - max(a['y'], c['y'])
            if ox > 0.04 and oy > 0.04:
                k = tuple(sorted([a['n'], c['n']]))
                if k in vus: continue
                vus.add(k)
                msgs.append(f"  CHEVAUCH {a['n']} x {c['n']} {ox:.2f}x{oy:.2f}\" «{a['t'][:22]}»/«{c['t'][:22]}»")

    if msgs:
        total += len(msgs)
        print(f"\n--- Diapositive {idx} ---")
        print("\n".join(msgs))

print(f"\n{total} signalement(s)." if total else "\nAucun signalement.")
